from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from typing import Any, Dict, List, Optional, Union

from pptx import Presentation


@dataclass
class TemplateInfo:
    slide_layout_names: List[str]
    masters: List[str]
    placeholder_map: Dict[str, Any]


def analyze_pptx_template(path_or_bytes: Union[Path, bytes]) -> TemplateInfo:
    if isinstance(path_or_bytes, bytes):
        prs = Presentation(BytesIO(path_or_bytes))
    else:
        prs = Presentation(str(path_or_bytes))
    layout_names = []
    for layout in prs.slide_layouts:
        name = getattr(layout, "name", None) or "layout"
        layout_names.append(name)

    masters = []
    for m in prs.slide_masters:
        masters.append(getattr(m, "name", "master"))

    # Best-effort placeholder map (layout -> placeholders)
    placeholder_map: Dict[str, Any] = {}
    for layout in prs.slide_layouts:
        lname = getattr(layout, "name", "layout")
        phs = []
        for shape in layout.placeholders:
            phs.append(
                {
                    "idx": shape.placeholder_format.idx,
                    "type": str(shape.placeholder_format.type),
                    "name": shape.name,
                }
            )
        placeholder_map[lname] = phs

    return TemplateInfo(
        slide_layout_names=layout_names,
        masters=masters,
        placeholder_map=placeholder_map,
    )
