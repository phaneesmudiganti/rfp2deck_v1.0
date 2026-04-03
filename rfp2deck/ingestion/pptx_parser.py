from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation


@dataclass
class ParsedPptx:
    text: str
    slide_count: int


def parse_pptx(path: Path) -> ParsedPptx:
    presentation = Presentation(path)
    texts: list[str] = []
    for i, slide in enumerate(presentation.slides, start=1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text_frame = shape.text_frame
            if not text_frame:
                continue
            for p in text_frame.paragraphs:
                t = (p.text or "").strip()
                if t:
                    slide_texts.append(t)
        if slide_texts:
            texts.append(f"\n\n--- SLIDE {i} ---\n\n" + "\n".join(slide_texts))
    return ParsedPptx(text="\n".join(texts).strip(), slide_count=len(presentation.slides))
