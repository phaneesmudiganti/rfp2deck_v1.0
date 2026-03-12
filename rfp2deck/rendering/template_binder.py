from __future__ import annotations
from dataclasses import dataclass
from typing import Dict, Optional, List
from pptx import Presentation


@dataclass
class BoundLayout:
    name: str
    title_ph_idx: Optional[int]
    body_ph_idx: Optional[int]
    picture_ph_idx: Optional[int]


class TemplateBinder:
    """Bind slide archetypes to real template layouts and placeholders."""

    def __init__(self, prs: Presentation):
        self.prs = prs
        self.layouts = list(prs.slide_layouts)
        self.bound: Dict[str, BoundLayout] = {}
        self._analyze_layouts()

    def _analyze_layouts(self) -> None:
        for layout in self.layouts:
            lname = getattr(layout, "name", "Layout")
            title_idx = None
            body_idx = None
            pic_idx = None
            for ph in layout.placeholders:
                pht = ph.placeholder_format.type
                if title_idx is None and "TITLE" in str(pht):
                    title_idx = ph.placeholder_format.idx
                if body_idx is None and ("BODY" in str(pht) or "CONTENT" in str(pht)):
                    body_idx = ph.placeholder_format.idx
                if pic_idx is None and ("PICTURE" in str(pht) or "IMAGE" in str(pht)):
                    pic_idx = ph.placeholder_format.idx
            self.bound[lname] = BoundLayout(lname, title_idx, body_idx, pic_idx)

    def list_layout_names(self) -> List[str]:
        return [getattr(l, "name", "Layout") for l in self.layouts]

    def pick_layout(self, archetype: str, layout_hint: Optional[str] = None):
        if layout_hint:
            for layout in self.layouts:
                if getattr(layout, "name", None) == layout_hint:
                    return layout

        wanted: List[str] = []
        a = (archetype or "").lower()
        if a == "title":
            wanted = ["title slide", "title"]
        elif a in (
            "agenda",
            "requirements",
            "customer context",
            "solution overview",
            "delivery plan",
            "risks",
            "case studies",
            "next steps",
            "content",
            "commercials",
        ):
            wanted = ["title and content", "title & content", "content", "section"]
        elif a == "team":
            wanted = ["two content", "comparison", "title and content", "content"]
        elif a == "timeline":
            wanted = ["timeline", "process", "title and content", "content"]
        elif a == "architecture":
            wanted = ["picture with caption", "picture", "title and content", "content"]
        else:
            wanted = ["title and content", "content"]

        for w in wanted:
            for layout in self.layouts:
                name = (getattr(layout, "name", "") or "").lower()
                if w in name:
                    return layout
        return self.layouts[0]

    def placeholders_for_layout(self, layout_name: str) -> BoundLayout:
        return self.bound.get(layout_name) or BoundLayout(layout_name, None, None, None)
