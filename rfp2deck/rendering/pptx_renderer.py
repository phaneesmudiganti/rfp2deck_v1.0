from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from rfp2deck.core.schemas import DeckPlan, DiagramSpec

# Layout constants (legacy defaults)
LEFT_MARGIN_IN = 0.75
RIGHT_MARGIN_IN = 0.75
TOP_MARGIN_IN = 0.45
BOTTOM_MARGIN_IN = 0.45

TITLE_X_IN = LEFT_MARGIN_IN
TITLE_Y_IN = TOP_MARGIN_IN
TITLE_H_IN = 0.70

CONTENT_TOP_Y_IN = 1.25
DIAGRAM_Y_IN = 1.35
DIAGRAM_H_IN = 3.55

BULLETS_Y_IN = 5.05
BULLETS_H_IN = 2.10

# Fonts
FONT_TITLE_PT = 34
FONT_BODY_START_PT = 18
FONT_BODY_MIN_PT = 12

EMU_PER_INCH = 914400


def _layout(prs: Presentation) -> dict:
    """Compute non-overlapping layout boxes sized from the actual slide size.

    This prevents text/image overlap when templates deviate slightly from 16:9.
    We keep a consistent hierarchy: title at top, large diagram in middle,
    bullets at bottom.
    """
    w_in = float(prs.slide_width) / EMU_PER_INCH
    h_in = float(prs.slide_height) / EMU_PER_INCH

    lm = max(0.55, w_in * 0.06)
    rm = lm
    tm = max(0.35, h_in * 0.06)
    bm = max(0.35, h_in * 0.06)

    title_x = lm
    title_y = tm
    title_w = w_in - lm - rm
    title_h = max(0.65, h_in * 0.11)

    diag_x = lm
    diag_y = title_y + title_h + (h_in * 0.03)
    diag_w = title_w
    diag_h = h_in * 0.50  # bigger, consistent diagrams

    bullets_x = lm
    bullets_y = diag_y + diag_h + (h_in * 0.03)
    bullets_w = title_w
    bullets_h = max(h_in - bullets_y - bm, h_in * 0.18)

    return {
        "title": (title_x, title_y, title_w, title_h),
        "diagram": (diag_x, diag_y, diag_w, diag_h),
        "bullets": (bullets_x, bullets_y, bullets_w, bullets_h),
    }


def _find_blank_layout(prs: Presentation):
    """Pick a layout with minimal placeholders."""
    # Prefer a truly blank layout if present
    for layout in prs.slide_layouts:
        if len(getattr(layout, "placeholders", [])) == 0:
            return layout
    # Otherwise pick the one with the fewest placeholders
    best = None
    best_n = 10**9
    for layout in prs.slide_layouts:
        n = len(getattr(layout, "placeholders", []))
        if n < best_n:
            best = layout
            best_n = n
    return best or prs.slide_layouts[0]


def _clear_text_on_slide(slide):
    """Clear all text frames on a slide (for template-clean rendering)."""
    for shape in slide.shapes:
        try:
            if getattr(shape, "has_text_frame", False):
                shape.text_frame.clear()
        except Exception:
            continue


def _remove_marker_shapes(slide):
    """Remove/clear template marker artifacts so they don't show in edit mode.

    We remove shapes that look like template guidance blocks ({{...}}) or
    PowerPoint placeholder guidance like 'Click to add...'.
    """
    to_remove = []
    for shape in slide.shapes:
        try:
            if getattr(shape, "has_text_frame", False):
                txt = (shape.text_frame.text or "").strip()
                upper = txt.upper()
                if (
                    ("{{" in txt and "}}" in txt)
                    or upper.startswith("TEMPLATE:")
                    or upper.startswith("CLICK TO ADD")
                    or upper in {"TITLE", "SUBTITLE", "CONTENT"}
                ):
                    to_remove.append(shape)
        except Exception:
            continue

    for shape in to_remove:
        try:
            slide.shapes._spTree.remove(shape._element)  # pylint: disable=protected-access
        except Exception:
            try:
                shape.text_frame.clear()
            except Exception:
                pass


def _add_title(
    slide,
    prs: Presentation,
    title: str,
    x_in: float | None = None,
    y_in: float | None = None,
    w_in: float | None = None,
    h_in: float | None = None,
):
    """Add a title box at the top of the slide.

    If coordinates are not provided, we fall back to the legacy constants.
    """
    if x_in is None or y_in is None or w_in is None or h_in is None:
        w = prs.slide_width / 914400.0
        box_w = w - LEFT_MARGIN_IN - RIGHT_MARGIN_IN
        x_in = TITLE_X_IN
        y_in = TITLE_Y_IN
        w_in = box_w
        h_in = TITLE_H_IN

    tb = slide.shapes.add_textbox(Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(FONT_TITLE_PT)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    tf.word_wrap = True


def _fit_font_for_box(lines, w_in: float, h_in: float) -> int:
    """Compute a readable font size to avoid overflow."""
    # very rough heuristic: more lines => smaller font
    n = max(1, len(lines))
    # scale by available vertical space
    base = FONT_BODY_START_PT
    # approximate: each bullet paragraph consumes ~0.26 inches at 18pt
    est_per_line = 0.26 * (base / 18.0)
    if n * est_per_line <= h_in:
        return base
    # shrink
    size = int(base * (h_in / max(n * est_per_line, 0.001)))
    return max(FONT_BODY_MIN_PT, min(base, size))


def _add_bullets(slide, x_in: float, y_in: float, w_in: float, h_in: float, bullets, font_pt: int):
    """Add bullets in a bounded box."""
    tb = slide.shapes.add_textbox(Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True

    first = True
    for b in bullets:
        txt = (b or "").strip()
        if not txt:
            continue
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = txt
        p.level = 0
        p.font.size = Pt(font_pt)
        p.alignment = PP_ALIGN.LEFT


def _place_image_contain(slide, img_path: Path, x_in: float, y_in: float, w_in: float, h_in: float):
    """Place an image contained within the bounding box (no overflow)."""
    from PIL import Image

    img = Image.open(img_path)
    iw, ih = img.size
    img.close()

    box_ratio = w_in / max(h_in, 1e-6)
    img_ratio = iw / max(ih, 1e-6)

    if img_ratio > box_ratio:
        # wider: fit width
        new_w = w_in
        new_h = w_in / img_ratio
    else:
        # taller: fit height
        new_h = h_in
        new_w = h_in * img_ratio

    px = x_in + (w_in - new_w) / 2.0
    py = y_in + (h_in - new_h) / 2.0

    slide.shapes.add_picture(
        str(img_path),
        Inches(px),
        Inches(py),
        width=Inches(new_w),
        height=Inches(new_h),
    )


def _render_consulting_slide(slide, prs: Presentation, title: str, bullets, diagram=None):
    """Stacked layout: title top, big image middle (if approved), bullets bottom."""
    # clear existing template artifacts
    _clear_text_on_slide(slide)
    _remove_marker_shapes(slide)

    layout = _layout(prs)
    tx, ty, tw, th = layout["title"]
    dx, dy, dw, dh = layout["diagram"]
    bx, by, bw, bh = layout["bullets"]

    _add_title(slide, prs, title, x_in=tx, y_in=ty, w_in=tw, h_in=th)

    # diagram if approved
    has_diagram = bool(
        diagram
        and getattr(diagram, "approved", False)
        and getattr(diagram, "image_path", None)
    )
    image_path = getattr(diagram, "image_path", None) if diagram else None

    if has_diagram:
        img = Path(str(image_path))
        if img.exists():
            _place_image_contain(slide, img, dx, dy, dw, dh)

        # bottom bullets
        font_pt = _fit_font_for_box(bullets, bw, bh)
        _add_bullets(slide, bx, by, bw, bh, bullets, font_pt)
        return

    # no diagram => use combined diagram+bullets region for text (still stacked)
    body_y = dy
    body_h = (by + bh) - dy
    font_pt = _fit_font_for_box(bullets, dw, body_h)
    _add_bullets(slide, dx, body_y, dw, body_h, bullets, font_pt)


def render_deck_from_template(deck_plan: DeckPlan, template_pptx: Path, out_path: Path) -> Path:
    """Render a new PPTX from a template and a deck plan."""
    prs = Presentation(str(template_pptx))
    blank_layout = _find_blank_layout(prs)

    # remove all existing slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId  # pylint: disable=protected-access
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]  # pylint: disable=protected-access

    for slide_spec in deck_plan.slides:
        slide = prs.slides.add_slide(blank_layout)

        _render_consulting_slide(
            slide,
            prs,
            slide_spec.title or "",
            slide_spec.bullets or [],
            diagram=getattr(slide_spec, "diagram", None),
        )

        # Final cleanup: remove any leftover template guidance/markers so edit mode stays clean.
        _remove_marker_shapes(slide)

    prs.save(str(out_path))
    return out_path